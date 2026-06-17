"""个人助手工具集（M5）。见 docs/design.md §4.6。

文件问答（read_dir / search_files / summarize_files，只读）+ 办公文档（xlsx/docx/pptx，写）+ draft。
所有文件操作都过 FolderAccess 授权校验（§5.2）：只能碰用户显式授权的目录，写操作走审批。
办公依赖（python-docx/pptx/openpyxl）惰性 import（pip install -e ".[office]"）。
"""
from __future__ import annotations

from typing import Any, Callable

from pydantic import BaseModel, Field

from core.harness.fs_access import FolderAccess
from core.tools.base import BaseTool, ToolContext, ToolResult, ValidationResult


class _FsTool(BaseTool):
    """带 FolderAccess 的工具基类：validate_input 先做授权校验（写工具早拦截，免无谓审批）。"""

    def __init__(self, fs: FolderAccess) -> None:
        self.fs = fs

    def _paths(self, inp: Any) -> list[str]:
        return [inp.path] if hasattr(inp, "path") else list(getattr(inp, "paths", []))

    async def validate_input(self, inp: Any, ctx: ToolContext) -> ValidationResult:
        for p in self._paths(inp):
            try:
                self.fs.resolve(p, need_write=not self.is_read_only(inp))
            except (PermissionError, FileNotFoundError) as e:
                return ValidationResult(ok=False, message=str(e))
        return ValidationResult(ok=True)


# ---------- 文件问答（只读）----------
class ReadDirInput(BaseModel):
    path: str = Field(description="授权目录内的文件夹路径")


class ReadDirTool(_FsTool):
    """列出授权目录内某文件夹的内容（文件与子目录）。"""

    name = "read_dir"
    input_model = ReadDirInput

    def is_read_only(self, inp: ReadDirInput) -> bool:
        return True

    async def call(self, inp: ReadDirInput, ctx: ToolContext, on_progress: Callable) -> ToolResult:
        d = self.fs.resolve(inp.path, must_exist=True)
        if not d.is_dir():
            raise NotADirectoryError(inp.path)
        lines = []
        for p in sorted(d.iterdir()):
            if p.is_dir():
                lines.append(f"📁 {p.name}/")
            else:
                lines.append(f"📄 {p.name} ({p.stat().st_size}B)")
        return ToolResult(data=f"{d}（{len(lines)} 项）:\n" + "\n".join(lines))


class SearchFilesInput(BaseModel):
    query: str = Field(description="要搜索的文本")
    path: str = Field(description="授权目录内搜索起点")
    max_results: int = Field(default=50, ge=1, le=500)


class SearchFilesTool(_FsTool):
    """在授权目录内按内容搜索文件（类似 grep），返回 文件:行号: 命中内容。"""

    name = "search_files"
    input_model = SearchFilesInput

    def is_read_only(self, inp: SearchFilesInput) -> bool:
        return True

    async def call(self, inp: SearchFilesInput, ctx: ToolContext, on_progress: Callable) -> ToolResult:
        base = self.fs.resolve(inp.path, must_exist=True)
        root = base if base.is_dir() else base.parent
        hits: list[str] = []
        for p in sorted(root.rglob("*")):
            if len(hits) >= inp.max_results:
                break
            if not p.is_file():
                continue
            try:
                text = p.read_text("utf-8", errors="ignore")
            except OSError:
                continue
            for i, line in enumerate(text.splitlines(), 1):
                if inp.query in line:
                    hits.append(f"{p}:{i}: {line.strip()[:120]}")
                    if len(hits) >= inp.max_results:
                        break
        return ToolResult(data="\n".join(hits) if hits else f"未找到 '{inp.query}'")


class SummarizeFilesInput(BaseModel):
    paths: list[str] = Field(description="授权目录内的文件路径列表")
    max_chars_each: int = Field(default=4000, ge=200, le=20000)


class SummarizeFilesTool(_FsTool):
    """读取多个授权目录内文件的内容，供你阅读/总结/问答。"""

    name = "summarize_files"
    input_model = SummarizeFilesInput

    def is_read_only(self, inp: SummarizeFilesInput) -> bool:
        return True

    async def call(self, inp: SummarizeFilesInput, ctx: ToolContext, on_progress: Callable) -> ToolResult:
        parts: list[str] = []
        for path in inp.paths[:20]:
            try:
                p = self.fs.resolve(path, must_exist=True)
                text = p.read_text("utf-8", errors="ignore")[: inp.max_chars_each]
                parts.append(f"=== {path} ===\n{text}")
            except Exception as e:  # noqa: BLE001
                parts.append(f"=== {path} ===\n[读取失败: {e}]")
        return ToolResult(data="\n\n".join(parts))


# ---------- 起草（纯文本，无副作用）----------
class DraftInput(BaseModel):
    kind: str = Field(description="草稿类型，如 邮件/周报/文档")
    content: str = Field(description="草稿正文或要点")


class DraftTool(BaseTool):
    """起草文本（邮件/周报/文档等）。纯文本产出，你审阅后自行使用，不落盘。"""

    name = "draft"
    input_model = DraftInput

    def is_read_only(self, inp: DraftInput) -> bool:
        return True

    async def call(self, inp: DraftInput, ctx: ToolContext, on_progress: Callable) -> ToolResult:
        return ToolResult(data=f"[草稿 · {inp.kind}]\n{inp.content}")


# ---------- 办公文档（写，走审批；惰性 import）----------
class XlsxWriteInput(BaseModel):
    path: str = Field(description="输出 .xlsx 路径（授权目录内）")
    rows: list[list] = Field(description="二维数组，每个子数组是一行（首行可作表头）")
    sheet: str = Field(default="Sheet1")


class XlsxWriteTool(_FsTool):
    """生成 Excel(.xlsx)：把二维数组 rows 写入工作表（授权目录内）。"""

    name = "xlsx_write"
    input_model = XlsxWriteInput

    async def call(self, inp: XlsxWriteInput, ctx: ToolContext, on_progress: Callable) -> ToolResult:
        from openpyxl import Workbook

        p = self.fs.resolve(inp.path, need_write=True)
        p.parent.mkdir(parents=True, exist_ok=True)
        wb = Workbook()
        ws = wb.active
        ws.title = inp.sheet
        for row in inp.rows:
            ws.append(list(row))
        wb.save(str(p))
        return ToolResult(data=f"已写入 {p}（{len(inp.rows)} 行）")


class DocxWriteInput(BaseModel):
    path: str = Field(description="输出 .docx 路径（授权目录内）")
    title: str | None = None
    paragraphs: list[str] = Field(default_factory=list)


class DocxWriteTool(_FsTool):
    """生成 Word(.docx)：可选标题 + 段落列表（授权目录内）。"""

    name = "docx_write"
    input_model = DocxWriteInput

    async def call(self, inp: DocxWriteInput, ctx: ToolContext, on_progress: Callable) -> ToolResult:
        from docx import Document

        p = self.fs.resolve(inp.path, need_write=True)
        p.parent.mkdir(parents=True, exist_ok=True)
        doc = Document()
        if inp.title:
            doc.add_heading(inp.title, level=1)
        for para in inp.paragraphs:
            doc.add_paragraph(para)
        doc.save(str(p))
        return ToolResult(data=f"已生成 {p}（标题{'有' if inp.title else '无'}，{len(inp.paragraphs)} 段）")


class PptxBuildInput(BaseModel):
    path: str = Field(description="输出 .pptx 路径（授权目录内）")
    slides: list[dict] = Field(description='每页 {"title": str, "bullets": [str,...]}')


class PptxBuildTool(_FsTool):
    """生成 PPT(.pptx)：slides 每页含 title 与 bullets 列表（授权目录内）。"""

    name = "pptx_build"
    input_model = PptxBuildInput

    async def call(self, inp: PptxBuildInput, ctx: ToolContext, on_progress: Callable) -> ToolResult:
        from pptx import Presentation

        p = self.fs.resolve(inp.path, need_write=True)
        p.parent.mkdir(parents=True, exist_ok=True)
        prs = Presentation()
        for s in inp.slides:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = str(s.get("title", ""))
            bullets = s.get("bullets", []) or []
            tf = slide.placeholders[1].text_frame
            if bullets:
                tf.text = str(bullets[0])
                for b in bullets[1:]:
                    tf.add_paragraph().text = str(b)
        prs.save(str(p))
        return ToolResult(data=f"已生成 {p}（{len(inp.slides)} 页）")
